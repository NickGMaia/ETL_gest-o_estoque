import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

df = pd.read_csv('data/movimentacao_estoque.csv')
df['Valor_Total_Venda'] = df['Quantidade'] * df['Valor_Unitario']
filiais = df['ID_Empresa'].unique()
print(filiais)
print("-" * 60)

df_saidas = df[df['Tipo'] == 'Saída'].copy()
df_entradas = df[df['Tipo'] == 'Entrada'].copy()

df_saidas['Faturamento'] = df_saidas['Quantidade'] * df_saidas['Valor_Unitario']

df['Variacao'] = np.where(df['Tipo'] == 'Entrada', 
                          df['Quantidade'], 
                          -df['Quantidade'])


total = df_saidas['Faturamento'].sum()
print(f"Total faturado: R$ {total:.2f}\n")
print("-" * 60)



df_vendas = df[df['Tipo'] == 'Saída'].copy()
relatorio_geral = df_vendas.groupby(['Categoria', 'Produto'])['Valor_Total_Venda'].sum().reset_index()

print(f'{relatorio_geral}\n')
print("-" * 60)

vendas_filial = df_vendas.groupby(['ID_Empresa', 'Produto'])['Valor_Total_Venda'].sum().reset_index()
print("Vendas por filial:")
print(f'{vendas_filial}\n')
print("-" * 60)
print("\nEstoque Atual:")

estoque_atual = df.groupby(['ID_Empresa', 'Produto'])['Variacao'].sum().reset_index()
estoque_atual.rename(columns={'Variacao': 'Saldo_Atual'}, inplace=True)
print(f'{estoque_atual}')
print("-" * 60)

erros_estoque = estoque_atual[estoque_atual['Saldo_Atual'] < 0]
print("\nATENÇÃO: Itens com erro de registro (Saldo Negativo):")
print(erros_estoque)
print("-" * 60)

Cat_prod_vendido = df_vendas.groupby(['ID_Empresa', 'Categoria'])['Quantidade'].sum().reset_index()
print('\nVenda por Filial:')
print (Cat_prod_vendido.sort_values(by='Quantidade', ascending=False))

estoque_atual['Status'] = np.where(estoque_atual['Saldo_Atual']<= 10,'⚠️ COMPRAR', '✅ OK')
print("\nRelatório de Reposição de Estoque:")
print(estoque_atual[estoque_atual['Status'] == '⚠️ COMPRAR'])
print("-" * 60)


ticket_medio = df_vendas.groupby('ID_Empresa')['Valor_Total_Venda'].mean().reset_index()
ticket_medio.rename(columns={'Valor_Total_Venda': 'Ticket_Medio'}, inplace=True)
print("\nTicket Médio por Operação:")
print(ticket_medio.sort_values(by='Ticket_Medio', ascending=False))
print("-" * 60)


df_vendas['Data'] = pd.to_datetime(df_vendas['Data'])
vendas_diarias = df_vendas.groupby(df_vendas['Data'].dt.date)['Valor_Total_Venda'].sum().reset_index()
print("\nEvolução de Vendas Diárias:")
print(vendas_diarias)
print("-" * 60)



# ===========================================================
# GERAÇÃO DO GRÁFICO 1: (LOAD)
# ===========================================================

plt.figure(figsize=(8, 5))

plt.bar(ticket_medio['ID_Empresa'], ticket_medio['Ticket_Medio'], color=['#38bdf8', '#0ea5e9', '#0284c7'])

plt.title('Ticket Médio por Operação em Cada Filial', fontsize=14, fontweight='bold')
plt.xlabel('Filiais', fontsize=12)
plt.ylabel('Valor Médio (R$)', fontsize=12)
plt.grid(axis='y', linestyle='--', alpha=0.7) # Linhas de grade ao fundo

plt.savefig('output/graficos/grafico_ticket_medio.png', bbox_inches='tight')
plt.close() # Fecha o gráfico para liberar a memória do computador



# ===========================================================
# GRÁFICO 2: EVOLUÇÃO DE VENDAS DIÁRIAS (Linha)
# ===========================================================

plt.figure(figsize=(12, 5))

plt.plot(vendas_diarias['Data'], vendas_diarias['Valor_Total_Venda'], 
         marker='o', color='#0ea5e9', linewidth=2)
plt.xticks(rotation=45, ha='right')
plt.title('Evolução do Faturamento Diário - Abril 2026', fontsize=14, fontweight='bold', pad=15)
plt.xlabel('Dias do Mês', fontsize=11)
plt.ylabel('Faturamento (R$)', fontsize=11)
plt.grid(True, linestyle='--', alpha=0.5)

plt.savefig('output/graficos/vendas_diarias.png', bbox_inches='tight')

plt.close()


# ===========================================================
# GRÁFICO 3: PARTICIPAÇÃO POR CATEGORIA (Pizza)
# ===========================================================

plt.figure(figsize=(6, 6))

vendas_categoria = df_vendas.groupby('Categoria')['Valor_Total_Venda'].sum().reset_index()

plt.pie(vendas_categoria['Valor_Total_Venda'], 
        labels=vendas_categoria['Categoria'], 
        autopct='%1.1f%%', 
        colors=['#38bdf8', '#34d399', '#f43f5e'], 
        startangle=140)

plt.title('Participação das Categorias no Faturamento', fontsize=14, fontweight='bold', pad=15)

plt.savefig('output/graficos/participacao_categorias.png', bbox_inches='tight')

plt.close()

print("[SUCESSO] Todos os gráficos foram gerados separadamente!")

# ===========================================================
# Converter para PPT
# ===========================================================

prs = Presentation()

slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Relatório de Gestão de Estoque"
subtitle.text = "Análise Automatizada de Performance - Abril 2026"

slide_layout = prs.slide_layouts[5] 
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Faturamento por Filial"
img_path = 'output/graficos/grafico_ticket_medio.png' 
slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))


slide_layout = prs.slide_layouts[5] 
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Participação por venda"
img_path = 'output/graficos/participacao_categorias.png' 
slide.shapes.add_picture(img_path, Inches(2.5), Inches(2), width=Inches(5))

slide_layout = prs.slide_layouts[5] 
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Vendas Diarias"
img_path = 'output/graficos/vendas_diarias.png' 
slide.shapes.add_picture(img_path, Inches(0), Inches(2), width=Inches(9.5))


slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Alertas de Estoque Crítico"

df_critico = estoque_atual[estoque_atual['Status'] == '⚠️ COMPRAR']

rows = len(df_critico) + 1
cols = 3
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(0.8)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.cell(0, 0).text = 'Filial'
table.cell(0, 1).text = 'Produto'
table.cell(0, 2).text = 'Saldo'

for i, (index, row) in enumerate(df_critico.iterrows()):
    table.cell(i+1, 0).text = str(row['ID_Empresa'])
    table.cell(i+1, 1).text = str(row['Produto'])
    table.cell(i+1, 2).text = str(row['Saldo_Atual'])


nome_arquivo = "Relatorio_Executivo_Estoque.pptx"
prs.save('output/Relatorio_Executivo_Estoque.pptx')

print(f"\n[FINALIZADO] O arquivo '{nome_arquivo}' foi criado com sucesso!")
